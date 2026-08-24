from __future__ import annotations

import base64
import csv
import html
import json
import mimetypes
import os
import re
import shutil
import subprocess
import sys
import tempfile
import urllib.error
import urllib.request
import zipfile
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import Callable, Iterable
from xml.etree import ElementTree


os.environ.setdefault("HF_HUB_DISABLE_SYMLINKS", "1")
os.environ.setdefault("HF_HUB_DISABLE_SYMLINKS_WARNING", "1")


SUPPORTED_EXTENSIONS = {
    ".bmp",
    ".csv",
    ".doc",
    ".docx",
    ".epub",
    ".jpeg",
    ".jpg",
    ".htm",
    ".html",
    ".json",
    ".log",
    ".md",
    ".odp",
    ".ods",
    ".odt",
    ".pdf",
    ".png",
    ".ppt",
    ".pptx",
    ".rtf",
    ".tif",
    ".tiff",
    ".txt",
    ".webp",
    ".xls",
    ".xlsx",
    ".xml",
}

DOCLING_EXTENSIONS = {
    ".bmp",
    ".csv",
    ".doc",
    ".docx",
    ".epub",
    ".htm",
    ".html",
    ".jpeg",
    ".jpg",
    ".md",
    ".odp",
    ".ods",
    ".odt",
    ".pdf",
    ".png",
    ".ppt",
    ".pptx",
    ".tif",
    ".tiff",
    ".webp",
    ".xls",
    ".xlsx",
}

LMSTUDIO_IMAGE_EXTENSIONS = {".bmp", ".jpeg", ".jpg", ".png", ".tif", ".tiff", ".webp"}
LMSTUDIO_EXTENSIONS = {".pdf", *LMSTUDIO_IMAGE_EXTENSIONS}
MINERU_EXTENSIONS = {".bmp", ".docx", ".jpeg", ".jpg", ".pdf", ".png", ".pptx", ".tif", ".tiff", ".webp", ".xlsx"}


@dataclass(frozen=True)
class ConversionResult:
    source: Path
    target: Path | None
    status: str
    message: str


@dataclass(frozen=True)
class ConversionOptions:
    markdown_engine: str = "docling"
    accelerator: str = "auto"
    pdf_backend: str = "docling_parse"
    ocr_engine: str = "rapidocr"
    ocr_mode: str = "full_page"
    quality: str = "high"
    ocr_scale: float | None = None
    ocr_lang: str = "de"
    rapidocr_text_score: float | None = None
    tesseract_psm: int | None = None
    easyocr_langs: tuple[str, ...] = ("de", "en")
    tesseract_langs: tuple[str, ...] = ("deu", "eng")
    table_mode: str = "accurate"
    table_structure_model: str = "v1"
    table_cell_matching: bool = True
    force_backend_text: bool = False
    layout_create_orphan_clusters: bool = True
    layout_keep_empty_clusters: bool = False
    layout_skip_cell_assignment: bool = False
    extract_pictures: bool = False
    describe_pictures: bool = False
    chart_extraction: bool = False
    heading_hierarchy: bool = True
    traverse_picture_text: bool = True
    escape_underscores: bool = False
    lmstudio_base_url: str = "http://localhost:1234/v1"
    lmstudio_model: str = ""
    lmstudio_max_tokens: int = 4096
    lmstudio_temperature: float = 0.0
    mineru_backend: str = "pipeline"
    mineru_method: str = "auto"
    mineru_lang: str = ""
    mineru_table: bool = True
    mineru_formula: bool = True
    mineru_image_analysis: bool = False
    mineru_api_url: str = ""


QUALITY_SETTINGS = {
    "fast": {"scale": 2.0, "picture_threshold": 0.05},
    "balanced": {"scale": 3.0, "picture_threshold": 0.03},
    "high": {"scale": 4.0, "picture_threshold": 0.015},
    "max": {"scale": 5.0, "picture_threshold": 0.005},
}


ProgressCallback = Callable[[ConversionResult], None]


def list_supported_files(input_dir: Path, recursive: bool = True, output_dir: Path | None = None) -> list[Path]:
    input_dir = input_dir.resolve()
    output_dir = output_dir.resolve() if output_dir else None
    pattern = "**/*" if recursive else "*"
    files: list[Path] = []

    for path in input_dir.glob(pattern):
        if not path.is_file():
            continue
        if output_dir and _is_relative_to(path.resolve(), output_dir):
            continue
        if path.suffix.lower() in SUPPORTED_EXTENSIONS:
            files.append(path)

    return sorted(files, key=lambda item: str(item).lower())


def convert_folder(
    input_dir: Path,
    output_dir: Path,
    recursive: bool = True,
    overwrite: bool = True,
    options: ConversionOptions | None = None,
    progress: ProgressCallback | None = None,
) -> list[ConversionResult]:
    input_dir = input_dir.resolve()
    output_dir = output_dir.resolve()

    if not input_dir.exists() or not input_dir.is_dir():
        raise ValueError(f"Input folder does not exist: {input_dir}")

    output_dir.mkdir(parents=True, exist_ok=True)

    results: list[ConversionResult] = []
    used_targets: set[Path] = set()

    for source in list_supported_files(input_dir, recursive=recursive, output_dir=output_dir):
        relative = source.relative_to(input_dir)
        target = _target_for(output_dir, relative, source.suffix.lower(), used_targets)
        used_targets.add(target)

        try:
            if target.exists() and not overwrite:
                result = ConversionResult(source, target, "skipped", "Markdown already exists")
            else:
                markdown = convert_file(source, options=options)
                target.parent.mkdir(parents=True, exist_ok=True)
                target.write_text(markdown, encoding="utf-8", newline="\n")
                result = ConversionResult(source, target, "created", "Markdown created")
        except MissingDependencyError as exc:
            result = ConversionResult(source, target, "skipped", str(exc))
        except EmptyExtractionError as exc:
            result = ConversionResult(source, target, "failed", str(exc))
        except Exception as exc:  # pragma: no cover - shown in the GUI log for user action.
            result = ConversionResult(source, target, "failed", f"{type(exc).__name__}: {exc}")

        results.append(result)
        if progress:
            progress(result)

    return results


def convert_file(source: Path, options: ConversionOptions | None = None) -> str:
    options = options or ConversionOptions()
    suffix = source.suffix.lower()

    body = ""
    docling_error: Exception | None = None

    if options.markdown_engine == "mineru" and suffix in MINERU_EXTENSIONS:
        body = _mineru_to_markdown(source, options)
        if not _has_content(body):
            raise EmptyExtractionError(f"MinerU did not return usable Markdown for {source.name}.")
    elif options.markdown_engine == "lmstudio" and suffix in LMSTUDIO_EXTENSIONS:
        body = _lmstudio_to_markdown(source, options)
        if not _has_content(body):
            raise EmptyExtractionError(f"LM Studio did not return usable Markdown for {source.name}.")
    elif suffix in DOCLING_EXTENSIONS:
        try:
            body = _docling_to_markdown(source, options)
        except MissingDependencyError as exc:
            docling_error = exc
        except Exception as exc:
            docling_error = exc

    if _has_content(body):
        pass
    elif suffix in {".txt", ".log", ".md"}:
        body = _read_text(source)
    elif suffix == ".csv":
        body = _csv_to_markdown(source)
    elif suffix == ".json":
        body = f"```json\n{json.dumps(json.loads(_read_text(source)), indent=2, ensure_ascii=False)}\n```"
    elif suffix in {".html", ".htm"}:
        body = _html_to_text(_read_text(source))
    elif suffix == ".xml":
        body = f"```xml\n{_read_text(source).strip()}\n```"
    elif suffix == ".rtf":
        body = _rtf_to_text(_read_text(source))
    elif suffix == ".pdf":
        body = _pdf_to_text(source)
    elif suffix == ".docx":
        body = _docx_to_markdown(source)
    elif suffix == ".xlsx":
        body = _xlsx_to_markdown(source)
    elif suffix == ".pptx":
        body = _pptx_to_markdown(source)
    elif suffix == ".odt":
        body = _odt_to_text(source)
    elif isinstance(docling_error, MissingDependencyError):
        raise docling_error
    elif docling_error:
        raise RuntimeError(f"Docling conversion failed: {docling_error}") from docling_error
    else:
        raise ValueError(f"Unsupported file type: {source.suffix}")

    if not _has_content(body):
        suffix_hint = " Scanned PDFs/images may need OCR support and a first run can take longer while models load."
        raise EmptyExtractionError(f"No text content could be extracted from {source.name}.{suffix_hint}")

    title = source.stem.replace("_", " ").replace("-", " ").strip() or source.name
    source_name = source.name.replace("\\", "/")
    return f"# {title}\n\n_Source: `{source_name}`_\n\n{body.strip()}\n"


class MissingDependencyError(RuntimeError):
    pass


class EmptyExtractionError(RuntimeError):
    pass


def _docling_to_markdown(path: Path, options: ConversionOptions) -> str:
    try:
        from docling.document_converter import DocumentConverter
    except ImportError as exc:
        raise MissingDependencyError("docling is missing; run install_windows.bat or pip install -r requirements.txt") from exc

    if path.stat().st_size == 0:
        raise EmptyExtractionError(f"{path.name} is empty and cannot be converted.")

    converter = _get_docling_converter(options)

    with tempfile.TemporaryDirectory(prefix="doclink_") as temp_dir:
        safe_path = Path(temp_dir) / f"source{path.suffix.lower()}"
        shutil.copy2(path, safe_path)
        result = converter.convert(safe_path)

    markdown = _export_docling_markdown(result.document, options)
    if options.describe_pictures:
        markdown = _append_picture_descriptions(markdown, result.document)

    return markdown


def _get_docling_converter(options: ConversionOptions):
    cache = getattr(_get_docling_converter, "_cache", {})
    if options not in cache:
        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.pipeline_options import PdfPipelineOptions
        from docling.document_converter import DocumentConverter
        from docling.document_converter import PdfFormatOption

        settings = _quality_settings(options)
        ocr_scale = _ocr_scale(options, settings)
        pipeline_options = PdfPipelineOptions()
        pipeline_options.accelerator_options = _accelerator_options(options)
        pipeline_options.do_ocr = options.ocr_engine != "none"
        pipeline_options.do_table_structure = options.table_mode != "off"
        pipeline_options.force_backend_text = options.force_backend_text
        pipeline_options.images_scale = ocr_scale
        pipeline_options.generate_picture_images = options.extract_pictures or options.describe_pictures
        pipeline_options.do_picture_description = options.describe_pictures
        pipeline_options.do_chart_extraction = options.chart_extraction
        _apply_layout_options(pipeline_options, options)

        if options.heading_hierarchy:
            try:
                from docling.datamodel.pipeline_options import HeadingHierarchyOptions

                pipeline_options.heading_hierarchy_options = HeadingHierarchyOptions(enabled=True)
            except Exception:
                pass

        if pipeline_options.do_table_structure:
            pipeline_options.table_structure_options = _get_table_structure_options(options)

        if options.ocr_engine != "none":
            pipeline_options.ocr_options = _get_ocr_options(options, ocr_scale)

        if options.describe_pictures:
            pipeline_options.picture_description_options = _get_picture_description_options(settings["picture_threshold"])

        pdf_format_kwargs: dict[str, object] = {"pipeline_options": pipeline_options}
        backend = _pdf_backend(options.pdf_backend)
        if backend is not None:
            pdf_format_kwargs["backend"] = backend

        cache[options] = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(**pdf_format_kwargs),
            }
        )
        _get_docling_converter._cache = cache  # type: ignore[attr-defined]
    return cache[options]


def _get_table_structure_options(options: ConversionOptions):
    from docling.datamodel.pipeline_options import TableFormerMode, TableStructureOptions

    if options.table_structure_model == "granite":
        try:
            from docling.datamodel.pipeline_options import GraniteVisionTableStructureOptions

            return GraniteVisionTableStructureOptions()
        except Exception:
            pass

    if options.table_structure_model == "v2":
        try:
            from docling.datamodel.pipeline_options import TableStructureV2Options

            return TableStructureV2Options(do_cell_matching=options.table_cell_matching)
        except Exception:
            pass

    table_options = TableStructureOptions(do_cell_matching=options.table_cell_matching)
    table_options.mode = TableFormerMode.FAST if options.table_mode == "fast" else TableFormerMode.ACCURATE
    return table_options


def _apply_layout_options(pipeline_options, options: ConversionOptions) -> None:
    layout_kwargs = {
        "create_orphan_clusters": options.layout_create_orphan_clusters,
        "keep_empty_clusters": options.layout_keep_empty_clusters,
        "skip_cell_assignment": options.layout_skip_cell_assignment,
    }

    try:
        from docling.datamodel.pipeline_options import LayoutObjectDetectionOptions

        pipeline_options.layout_options = LayoutObjectDetectionOptions(**layout_kwargs)
    except Exception:
        try:
            from docling.datamodel.pipeline_options import LayoutOptions

            pipeline_options.layout_options = LayoutOptions(**layout_kwargs)
        except Exception:
            pass

    try:
        from docling.datamodel.pipeline_options import LayoutPostprocessorOptions

        pipeline_options.layout_postprocessor_options = LayoutPostprocessorOptions(**layout_kwargs)
    except Exception:
        pass


def _pdf_backend(value: str):
    if value == "auto":
        return None

    imports = {
        "docling_parse": ("docling.backend.docling_parse_backend", "DoclingParseDocumentBackend"),
        "threaded_docling_parse": ("docling.backend.docling_parse_backend", "ThreadedDoclingParseDocumentBackend"),
        "pypdfium2": ("docling.backend.pypdfium2_backend", "PyPdfiumDocumentBackend"),
        "dlparse_v1": ("docling.backend.docling_parse_v1_backend", "DoclingParseV1DocumentBackend"),
        "dlparse_v2": ("docling.backend.docling_parse_v2_backend", "DoclingParseV2DocumentBackend"),
        "dlparse_v4": ("docling.backend.docling_parse_v4_backend", "DoclingParseV4DocumentBackend"),
    }
    module_name, class_name = imports.get(value, imports["docling_parse"])
    try:
        module = __import__(module_name, fromlist=[class_name])
        return getattr(module, class_name)
    except Exception:
        return None


def _get_ocr_options(options: ConversionOptions, scale: float):
    from docling.datamodel.pipeline_options import OcrMode

    mode = _ocr_mode(options.ocr_mode, OcrMode)

    if options.ocr_engine == "auto":
        try:
            from docling.datamodel.pipeline_options import OcrAutoOptions

            return OcrAutoOptions(mode=mode, scale=scale)
        except Exception:
            pass

    if options.ocr_engine == "tesseract_cli":
        from docling.datamodel.pipeline_options import TesseractCliOcrOptions

        _ensure_tesseract_available()
        kwargs: dict[str, object] = {"mode": mode, "lang": list(options.tesseract_langs), "scale": scale}
        if options.tesseract_psm is not None:
            kwargs["psm"] = options.tesseract_psm
        tesseract_cmd = os.getenv("TESSERACT_CMD")
        if tesseract_cmd:
            kwargs["tesseract_cmd"] = tesseract_cmd
        return TesseractCliOcrOptions(**kwargs)

    if options.ocr_engine == "easyocr":
        from docling.datamodel.pipeline_options import EasyOcrOptions

        try:
            import easyocr  # noqa: F401
        except ImportError as exc:
            raise MissingDependencyError("EasyOCR is missing; run install_windows.bat again.") from exc

        return EasyOcrOptions(mode=mode, lang=list(options.easyocr_langs), scale=scale)

    if options.ocr_engine == "rapidocr":
        from docling.datamodel.pipeline_options import RapidOcrOptions

        kwargs: dict[str, object] = {"mode": mode, "lang": [options.ocr_lang], "scale": scale}
        if options.rapidocr_text_score is not None:
            kwargs["text_score"] = options.rapidocr_text_score
        return RapidOcrOptions(**kwargs)

    raise MissingDependencyError(f"Unknown OCR engine: {options.ocr_engine}")


def _accelerator_options(options: ConversionOptions):
    from docling.datamodel.accelerator_options import AcceleratorDevice, AcceleratorOptions

    devices = {
        "auto": AcceleratorDevice.AUTO,
        "cpu": AcceleratorDevice.CPU,
        "cuda": AcceleratorDevice.CUDA,
    }
    return AcceleratorOptions(device=devices.get(options.accelerator, AcceleratorDevice.AUTO))


def _ensure_tesseract_available() -> None:
    tesseract_cmd = os.getenv("TESSERACT_CMD") or shutil.which("tesseract")
    if not tesseract_cmd:
        raise MissingDependencyError(
            "Tesseract was not found. Install Tesseract OCR for Windows and make sure tesseract.exe is in PATH, "
            "or set TESSERACT_CMD to the full path."
        )

    try:
        subprocess.run([tesseract_cmd, "--version"], capture_output=True, text=True, timeout=10, check=True)
    except Exception as exc:
        raise MissingDependencyError(f"Tesseract could not be started: {tesseract_cmd}") from exc


def _ocr_mode(value: str, enum_cls):
    mapping = {
        "default": enum_cls.DEFAULT,
        "full_page": enum_cls.FULL_PAGE,
        "layout_regions": enum_cls.LAYOUT_REGIONS,
        "pdf_aware_layout_regions": enum_cls.PDF_AWARE_LAYOUT_REGIONS,
    }
    return mapping.get(value, enum_cls.FULL_PAGE)


def _get_picture_description_options(picture_threshold: float):
    prompt = (
        "Transkribiere allen sichtbaren Text im Bild exakt. "
        "Wenn kein Text vorhanden ist, beschreibe das Bild kurz und sachlich auf Deutsch."
    )

    try:
        from docling.datamodel.pipeline_options import PictureDescriptionVlmEngineOptions

        picture_options = PictureDescriptionVlmEngineOptions.from_preset("smolvlm")
    except Exception as exc:
        raise MissingDependencyError(
            "Picture description needs Docling VLM support; run install_vision_windows.bat."
        ) from exc

    picture_options.prompt = prompt
    picture_options.picture_area_threshold = picture_threshold
    picture_options.scale = 2.0
    return picture_options


def _append_picture_descriptions(markdown: str, document) -> str:
    lines: list[str] = []

    for index, picture in enumerate(getattr(document, "pictures", []), start=1):
        meta = getattr(picture, "meta", None)
        description = getattr(meta, "description", None) if meta else None
        text = getattr(description, "text", "") if description else ""
        caption = ""
        try:
            caption = picture.caption_text(doc=document)
        except Exception:
            caption = ""

        if text.strip() or caption.strip():
            lines.append(f"### Bild {index}")
            if caption.strip():
                lines.append(f"Beschriftung: {caption.strip()}")
            if text.strip():
                lines.append(text.strip())
            lines.append("")

    if not lines:
        return markdown

    return f"{markdown.rstrip()}\n\n## Bildtexte und Beschreibungen\n\n" + "\n".join(lines).rstrip() + "\n"


def _export_docling_markdown(document, options: ConversionOptions) -> str:
    kwargs = {
        "escape_html": False,
        "escape_underscores": options.escape_underscores,
        "image_placeholder": "" if not options.extract_pictures else "<!-- image embedded -->",
        "include_annotations": True,
        "enable_chart_tables": True,
        "compact_tables": False,
        "traverse_pictures": options.traverse_picture_text or options.describe_pictures or options.ocr_mode == "full_page",
    }

    if options.extract_pictures:
        try:
            from docling_core.types.doc import ImageRefMode

            kwargs["image_mode"] = ImageRefMode.EMBEDDED
        except Exception:
            pass

    try:
        return document.export_to_markdown(**kwargs)
    except TypeError:
        return document.export_to_markdown()


def _quality_settings(options: ConversionOptions) -> dict[str, float]:
    return QUALITY_SETTINGS.get(options.quality, QUALITY_SETTINGS["balanced"])


def _ocr_scale(options: ConversionOptions, settings: dict[str, float]) -> float:
    if options.ocr_scale is not None and options.ocr_scale > 0:
        return options.ocr_scale
    return settings["scale"]


def _mineru_to_markdown(path: Path, options: ConversionOptions) -> str:
    mineru = _mineru_executable()
    with tempfile.TemporaryDirectory(prefix="doclink_mineru_") as temp_dir:
        output_dir = Path(temp_dir) / "output"
        command = [
            str(mineru),
            "-p",
            str(path),
            "-o",
            str(output_dir),
            "-b",
            options.mineru_backend,
            "-m",
            options.mineru_method,
            "-t",
            _bool_arg(options.mineru_table),
            "-f",
            _bool_arg(options.mineru_formula),
        ]

        if options.mineru_lang.strip():
            command.extend(["-l", options.mineru_lang.strip()])
        if options.mineru_image_analysis:
            command.extend(["--image-analysis", "true"])
        if options.mineru_api_url.strip():
            command.extend(["--api-url", options.mineru_api_url.strip()])

        env = os.environ.copy()
        env.setdefault("PYTHONUTF8", "1")
        env.setdefault("HF_HUB_DISABLE_SYMLINKS", "1")
        env.setdefault("HF_HUB_DISABLE_SYMLINKS_WARNING", "1")
        completed = subprocess.run(command, capture_output=True, text=True, encoding="utf-8", errors="replace", env=env)
        if completed.returncode != 0:
            details = "\n".join(part.strip() for part in (completed.stdout, completed.stderr) if part.strip())
            raise RuntimeError(f"MinerU failed with exit code {completed.returncode}: {details[-2000:]}")

        markdown_files = sorted(output_dir.rglob("*.md"), key=lambda item: item.stat().st_size, reverse=True)
        if not markdown_files:
            details = "\n".join(part.strip() for part in (completed.stdout, completed.stderr) if part.strip())
            raise EmptyExtractionError(f"MinerU created no Markdown for {path.name}. {details[-1000:]}")

        return markdown_files[0].read_text(encoding="utf-8", errors="replace")


def _mineru_executable() -> Path:
    executable_name = "mineru.exe" if sys.platform.startswith("win") else "mineru"
    local_candidate = Path(sys.executable).resolve().parent / executable_name
    if local_candidate.exists():
        return local_candidate

    found = shutil.which("mineru")
    if found:
        return Path(found)

    raise MissingDependencyError("MinerU is missing; run install_mineru_windows.bat first.")


def _bool_arg(value: bool) -> str:
    return "true" if value else "false"


def _lmstudio_to_markdown(path: Path, options: ConversionOptions) -> str:
    if path.suffix.lower() == ".pdf":
        images = _render_pdf_pages(path, options)
    elif path.suffix.lower() in LMSTUDIO_IMAGE_EXTENSIONS:
        images = [(1, path.read_bytes(), _mime_type(path))]
    else:
        raise ValueError(f"LM Studio mode supports PDFs and images, not {path.suffix}")

    model = options.lmstudio_model.strip() or _lmstudio_first_model(options.lmstudio_base_url)
    sections: list[str] = []
    total_pages = len(images)

    for page_number, image_bytes, mime_type in images:
        page_markdown = _lmstudio_page_to_markdown(
            image_bytes=image_bytes,
            mime_type=mime_type,
            page_number=page_number,
            total_pages=total_pages,
            model=model,
            options=options,
        )
        sections.append(page_markdown.strip())

    return "\n\n".join(section for section in sections if section.strip())


def _render_pdf_pages(path: Path, options: ConversionOptions) -> list[tuple[int, bytes, str]]:
    try:
        import pypdfium2 as pdfium
    except ImportError as exc:
        raise MissingDependencyError("pypdfium2 is missing; run install_windows.bat again.") from exc

    settings = _quality_settings(options)
    pages: list[tuple[int, bytes, str]] = []

    with tempfile.TemporaryDirectory(prefix="doclink_lmstudio_") as temp_dir:
        pdf = pdfium.PdfDocument(str(path))
        try:
            for index in range(len(pdf)):
                page = pdf[index]
                try:
                    bitmap = page.render(scale=settings["scale"])
                    image = bitmap.to_pil()
                    if image.mode not in {"RGB", "L"}:
                        image = image.convert("RGB")
                    image_path = Path(temp_dir) / f"page_{index + 1}.jpg"
                    image.save(image_path, format="JPEG", quality=92, optimize=True)
                    pages.append((index + 1, image_path.read_bytes(), "image/jpeg"))
                finally:
                    page.close()
        finally:
            pdf.close()

    return pages


def _lmstudio_page_to_markdown(
    image_bytes: bytes,
    mime_type: str,
    page_number: int,
    total_pages: int,
    model: str,
    options: ConversionOptions,
) -> str:
    data_url = f"data:{mime_type};base64,{base64.b64encode(image_bytes).decode('ascii')}"
    prompt = _lmstudio_strict_ocr_prompt(page_number, total_pages)
    payload = {
        "model": model,
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            }
        ],
        "temperature": options.lmstudio_temperature,
        "top_p": 0.05,
        "max_tokens": options.lmstudio_max_tokens,
        "stream": False,
    }
    response = _lmstudio_request(options.lmstudio_base_url, "/chat/completions", payload)
    try:
        content = response["choices"][0]["message"]["content"]
    except (KeyError, IndexError, TypeError) as exc:
        raise RuntimeError(f"Unexpected LM Studio response: {response}") from exc

    if isinstance(content, list):
        markdown = "\n".join(item.get("text", "") for item in content if isinstance(item, dict)).strip()
    else:
        markdown = str(content).strip()

    return _clean_lmstudio_markdown(markdown)


def _lmstudio_strict_ocr_prompt(page_number: int, total_pages: int) -> str:
    return (
        "You are a strict OCR transcription engine for scanned administrative documents.\n"
        "Task: convert exactly this single visible page to raw Markdown. Treat the page image as the only source of truth.\n"
        "\n"
        "Hard rules:\n"
        "- Transcribe only text that is visibly present on the page. Never infer, complete, calculate, correct, or invent text.\n"
        "- Keep the original language, spelling, punctuation, numbers, dates, currency values, line breaks, and reading order.\n"
        "- Do not translate. Do not summarize. Do not explain. Do not add comments, warnings, confidence notes, or placeholders.\n"
        "- If a word, number, or cell is not readable, leave that part empty. Do not write 'unreadable', 'illegible', or similar.\n"
        "- Do not repeat text to fill visually empty areas. Empty table cells must stay empty.\n"
        "\n"
        "Text size and formatting:\n"
        "- Preserve relative text size and hierarchy as closely as Markdown allows.\n"
        "- Use # only for the largest visible page title. Use ## for clearly smaller section headings. Use ### only for "
        "clear subheadings. Do not make ordinary body text into headings.\n"
        "- Use **bold** only when the source text is visibly bold or functions as a printed label. Do not add emphasis yourself.\n"
        "- Keep small footer/header text as normal text unless it is visibly a heading. Do not enlarge it.\n"
        "- Preserve visible indentation, lists, numbering, labels, and separate blocks.\n"
        "\n"
        "Tables and forms:\n"
        "- Reconstruct tables only from visible rows, columns, borders, labels, and cell positions.\n"
        "- Keep the same column order and row order as the page.\n"
        "- Leave empty cells empty. Never copy a label or description into neighboring columns such as 'von', 'bis', "
        "amount, date, or factor columns unless the text is visibly printed there.\n"
        "- If Markdown table syntax cannot represent the layout without moving text into wrong cells, use an HTML table "
        "with <table>, <tr>, <th>, and <td>; empty cells must be <td></td>.\n"
        "- Use colspan/rowspan only when a cell is visibly merged in the document. Otherwise keep separate empty cells.\n"
        "\n"
        "Output:\n"
        "- Return raw Markdown only. No fenced code block. No prose before or after.\n"
        f"- Page context for you only: page {page_number} of {total_pages}. Do not print this context unless it is visibly printed."
    )


def _clean_lmstudio_markdown(markdown: str) -> str:
    cleaned = markdown.strip()
    cleaned = re.sub(r"^\s*```(?:markdown|md)?\s*\n", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\n\s*```\s*$", "", cleaned)
    cleaned = re.sub(r"\n\s*```(?:markdown|md)?\s*\n", "\n\n", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\n\s*```\s*\n", "\n\n", cleaned)
    return cleaned.strip()


def _lmstudio_first_model(base_url: str) -> str:
    response = _lmstudio_request(base_url, "/models", None, method="GET")
    try:
        return str(response["data"][0]["id"])
    except (KeyError, IndexError, TypeError) as exc:
        raise MissingDependencyError("LM Studio has no loaded model. Load a vision model in LM Studio first.") from exc


def _lmstudio_request(base_url: str, endpoint: str, payload: dict | None, method: str = "POST") -> dict:
    url = _normalize_lmstudio_base_url(base_url) + endpoint
    data = json.dumps(payload).encode("utf-8") if payload is not None else None
    request = urllib.request.Request(
        url,
        data=data,
        method=method,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {os.getenv('LMSTUDIO_API_KEY', 'lm-studio')}",
        },
    )
    try:
        with urllib.request.urlopen(request, timeout=300) as response:
            return json.loads(response.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        details = exc.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"LM Studio HTTP {exc.code}: {details}") from exc
    except urllib.error.URLError as exc:
        raise MissingDependencyError(
            f"LM Studio server is not reachable at {base_url}. Start the local server in LM Studio."
        ) from exc


def _normalize_lmstudio_base_url(base_url: str) -> str:
    base = base_url.strip().rstrip("/") or "http://localhost:1234/v1"
    if not base.endswith("/v1"):
        base = f"{base}/v1"
    return base


def _mime_type(path: Path) -> str:
    return mimetypes.guess_type(path.name)[0] or "image/png"


def _has_content(markdown: str) -> bool:
    text = re.sub(r"!\[[^\]]*]\([^)]*\)", "", markdown)
    text = re.sub(r"\bpage\s+\d+\b", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\bslide\s+\d+\b", "", text, flags=re.IGNORECASE)
    text = re.sub(r"#+\s*", "", text)
    text = re.sub(r"[_`|\\-]", " ", text)
    return bool(re.search(r"[A-Za-z0-9À-ÖØ-öø-ÿ]{3,}", text))


def _target_for(output_dir: Path, relative: Path, suffix: str, used_targets: set[Path]) -> Path:
    target = output_dir / relative.with_suffix(".md")
    if target not in used_targets:
        return target

    alternate = output_dir / relative.with_name(f"{relative.stem}{suffix.replace('.', '_')}.md")
    counter = 2
    while alternate in used_targets:
        alternate = output_dir / relative.with_name(f"{relative.stem}{suffix.replace('.', '_')}_{counter}.md")
        counter += 1
    return alternate


def _read_text(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(errors="replace")


def _csv_to_markdown(path: Path) -> str:
    text = _read_text(path)
    sample = text[:2048]
    try:
        dialect = csv.Sniffer().sniff(sample) if sample.strip() else csv.excel
    except csv.Error:
        dialect = csv.excel

    rows = list(csv.reader(text.splitlines(), dialect))
    if not rows:
        return ""

    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    header = normalized[0]
    separator = ["---"] * width
    lines = [_markdown_row(header), _markdown_row(separator)]
    lines.extend(_markdown_row(row) for row in normalized[1:])
    return "\n".join(lines)


def _markdown_row(values: Iterable[str]) -> str:
    escaped = [str(value).replace("|", "\\|").replace("\n", " ").strip() for value in values]
    return "| " + " | ".join(escaped) + " |"


class _TextHTMLParser(HTMLParser):
    block_tags = {"br", "div", "h1", "h2", "h3", "h4", "h5", "h6", "li", "p", "td", "th", "tr"}

    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in self.block_tags:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self.block_tags:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if data.strip():
            self.parts.append(html.unescape(data.strip()))

    def text(self) -> str:
        return "\n".join(line.strip() for line in "".join(self.parts).splitlines() if line.strip())


def _html_to_text(content: str) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        parser = _TextHTMLParser()
        parser.feed(content)
        return parser.text()

    soup = BeautifulSoup(content, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text("\n", strip=True)


def _rtf_to_text(content: str) -> str:
    content = re.sub(r"\\'[0-9a-fA-F]{2}", " ", content)
    content = re.sub(r"\\[a-zA-Z]+\d* ?", " ", content)
    content = content.replace("{", " ").replace("}", " ")
    content = content.replace("\\", "")
    return "\n".join(line.strip() for line in content.splitlines() if line.strip())


def _pdf_to_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise MissingDependencyError("pypdf is missing; run install_windows.bat or pip install -r requirements.txt") from exc

    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"## Page {index}\n\n{text.strip()}")
    return "\n\n".join(pages)


def _docx_to_markdown(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise MissingDependencyError("python-docx is missing; run install_windows.bat or pip install -r requirements.txt") from exc

    document = Document(str(path))
    lines: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            style = paragraph.style.name.lower() if paragraph.style else ""
            if style.startswith("heading"):
                level_match = re.search(r"(\d+)", style)
                level = min(int(level_match.group(1)), 6) if level_match else 2
                lines.append(f"{'#' * level} {text}")
            else:
                lines.append(text)
            lines.append("")

    for table in document.tables:
        rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
        if rows:
            width = max(len(row) for row in rows)
            normalized = [row + [""] * (width - len(row)) for row in rows]
            lines.append(_markdown_row(normalized[0]))
            lines.append(_markdown_row(["---"] * width))
            lines.extend(_markdown_row(row) for row in normalized[1:])
            lines.append("")

    return "\n".join(lines)


def _xlsx_to_markdown(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise MissingDependencyError("openpyxl is missing; run install_windows.bat or pip install -r requirements.txt") from exc

    workbook = load_workbook(path, read_only=True, data_only=True)
    sections: list[str] = []

    for sheet in workbook.worksheets:
        sections.append(f"## {sheet.title}")
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                rows.append(values)
        if rows:
            width = max(len(row) for row in rows)
            normalized = [row + [""] * (width - len(row)) for row in rows]
            sections.append(_markdown_row(normalized[0]))
            sections.append(_markdown_row(["---"] * width))
            sections.extend(_markdown_row(row) for row in normalized[1:])
        sections.append("")

    return "\n".join(sections)


def _pptx_to_markdown(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise MissingDependencyError("python-pptx is missing; run install_windows.bat or pip install -r requirements.txt") from exc

    presentation = Presentation(str(path))
    sections: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        sections.append(f"## Slide {index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                sections.append(shape.text.strip())
        sections.append("")

    return "\n".join(sections)


def _odt_to_text(path: Path) -> str:
    namespaces = {
        "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
    }

    with zipfile.ZipFile(path) as archive:
        content = archive.read("content.xml")

    root = ElementTree.fromstring(content)
    lines: list[str] = []
    for paragraph in root.findall(".//text:p", namespaces):
        text = "".join(paragraph.itertext()).strip()
        if text:
            lines.append(text)
    return "\n\n".join(lines)


def _is_relative_to(path: Path, parent: Path) -> bool:
    try:
        path.relative_to(parent)
        return True
    except ValueError:
        return False
